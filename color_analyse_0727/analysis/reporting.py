"""Reporting: parameters JSON, README index, output_index.csv, report.md, PPTX."""

from __future__ import annotations

import csv
import json
import re
from datetime import date
from pathlib import Path
from typing import Any, Iterable

import numpy as np
import pandas as pd

from pipeline.spectral_features import FEATURE_BANDS, NOISE_BANDS
from analysis.common import (
    DEFAULT_WINDOW_LABELS,
    FS,
    RADIUS_MM,
    ROI_TERMS,
    SUBJECTS,
    AnalysisVariant,
    all_variants,
    df_to_markdown,
)


def write_parameters(
    out: Path,
    *,
    variants: Iterable[AnalysisVariant],
    perms: int,
    workers: int,
    nperseg: int,
    noverlap: int,
) -> Path:
    parameters = {
        "analysis": "color_analyse_0727 v2",
        "date": date.today().isoformat(),
        "subjects": list(SUBJECTS),
        "default_window": list(DEFAULT_WINDOW_LABELS),
        "windows": [list(v.window) for v in variants],
        "signals": list(dict.fromkeys(v.signal for v in variants)),
        "variants": [v.suffix for v in variants],
        "signal_bands_hz": {
            signal: (list(band) if band else None)
            for signal, band in {
                "lf30": (1.0, 30.0),
                "raw200": None,
            }.items()
        },
        "baseline_ms": [-200.0, 0.0],
        "feature_bands_hz": [list(b) for b in FEATURE_BANDS],
        "removed_line_noise_bands_hz": [list(b) for b in sorted(NOISE_BANDS)],
        "band_power_method": "Welch log power",
        "nperseg": nperseg,
        "noverlap": noverlap,
        "fs_hz": FS,
        "n_permutations": perms,
        "workers": workers,
        "radius_mm": RADIUS_MM,
        "roi_terms": list(ROI_TERMS),
        "fdr_policy": "no cross-channel or within-channel FDR gate; informational columns only",
        "strategy1_rule": "ANOVA two-way color x category, type-II color main effect p<0.05 (2026-08-05 standard)",
        "strategy1_merged_rule": "historical pooled color-vs-gray (all four categories merged) MWU p<0.05, informational",
        "strategy2_rule": "at least one Task-1 category reaches Welch t-test p<0.05",
    }
    path = out / "analysis_parameters.json"
    path.write_text(
        json.dumps(parameters, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    return path


_DESCRIPTIONS = (
    ("functional_selection", "功能筛选 S1/S2 (functional selection)"),
    ("electrode_sets_and_csc", "CSC 电极集合 (electrode sets and CSC)"),
    ("electrode_set_summary_by_subject", "电极集合计数 (electrode set counts)"),
    ("patch_coordinates", "Talairach-MNI patch 坐标 (patch coordinates)"),
    ("csc_amplitude_statistics", "幅度统计 (amplitude statistics)"),
    ("csc_spectral_band_statistics", "16 频带功率统计 (band-power statistics)"),
    ("csc_spectral_group_friedman", "subject 级 Friedman (group spectral)"),
    ("decoding_summary", "频谱级解码汇总 (spectrum decoding summary)"),
    ("decoding_direction_consistency", "记忆颜色类内方向一致性 (sign consistency)"),
    ("luminance", "刺激亮度审计 (stimulus luminance audit)"),
)


def _describe(path: Path) -> str:
    name = path.name
    for prefix, description in _DESCRIPTIONS:
        if name.startswith(prefix):
            return description
    return "分析输出 (analysis output)"


def _variant_of(path: Path) -> tuple[str, str]:
    match = re.search(r"_(\d+-\d+)_(lf|broadband)\.(?:csv|png)$", path.name)
    if match:
        return match.group(1), match.group(2)
    return "", ""


def write_output_index(out: Path) -> Path:
    rows: list[dict[str, str]] = []
    for stage in sorted(
        p for p in out.iterdir() if p.is_dir() and p.name.startswith("stage")
    ):
        for path in sorted(stage.rglob("*")):
            if not path.is_file():
                continue
            window, signal = _variant_of(path)
            subject_match = re.search(r"(test00\d)", path.name)
            rows.append(
                {
                    "path": str(path.relative_to(out)),
                    "stage": stage.name,
                    "window": window,
                    "signal": signal,
                    "subject": subject_match.group(1) if subject_match else "",
                    "description": _describe(path),
                }
            )
    path = out / "output_index.csv"
    with path.open("w", newline="", encoding="utf-8-sig") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)
    return path


def write_readme(out: Path, perms: int, workers: int) -> Path:
    lines = [
        "# 颜色分析管线 v2 结果目录",
        "",
        f"生成日期：{date.today().isoformat()}",
        "",
        "## 目录结构",
        "",
        "- `analysis_parameters.json`：本次运行全部参数",
        "- `output_index.csv`：全部输出文件的索引（可用 pandas 按 stage/window/signal/subject 过滤）",
        "- `report/`：最终报告 final_analysis_report.md 与 .pptx",
        "- `stage01_selection/`：功能筛选、N2 空间筛选、CSC 电极集合与重叠图",
        "- `stage02_amplitude_spectral/`：幅度统计、16 频带功率统计与频谱图",
        "- `stage03_decoding/`：频谱级解码、置换 p 值、类内方向一致性",
        "- `stage04_luminance/`：刺激亮度/对比度/色彩度审计",
        "- `cache/`：中间特征数组（不浏览）",
        "",
        "## 命名规则",
        "",
        "变体专属文件带 `{时间窗}_{信号}` 后缀（`0-300_lf`、`100-400_broadband` 等）；",
        "所有表都带 `window` 与 `signal` 列，便于过滤。",
        "",
        "## 重跑命令",
        "",
        "```powershell",
        f"C:\\Users\\saber_soul\\.conda\\envs\\lr0727\\python.exe -m analysis.run_final_analysis --perms {perms} --workers {workers}",
        "```",
        "",
        "只跑部分阶段：`--stages selection csc stats decoding luminance report` 中任选；",
        "当前主变体：`--windows 100-400 --signals lf30 raw200`；也可显式指定其他窗口做敏感性分析。",
        "",
        "## 浏览方式",
        "",
        "打开 `notebooks/02_browse_results.ipynb`，设置 RESULT_DIR 后可按阶段/窗口/信号/被试过滤查看。",
        "",
    ]
    path = out / "README.md"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def _variant_counts_table(
    variant_tables: dict[str, dict[str, Any]],
) -> pd.DataFrame:
    rows = []
    for suffix, tables in variant_tables.items():
        summary = tables.get("summary")
        functional = tables.get("functional")
        if summary is not None and not summary.empty:
            rows.append(
                {
                    "variant": suffix,
                    "common_centers": int(summary.common_all_task_centers.sum()),
                    "strategy1": int(summary.strategy1.sum()),
                    "strategy2": int(summary.strategy2.sum()),
                    "N2_union": int(summary.N2_union.sum()),
                    "CSC": int(summary.CSC.sum()),
                    "CSC_merged": int(
                        summary.CSC_merged.sum()
                        if "CSC_merged" in summary.columns
                        else 0
                    ),
                    "CSC_fdr": int(
                        summary.CSC_fdr.sum()
                        if "CSC_fdr" in summary.columns
                        else 0
                    ),
                }
            )
        elif functional is not None and not functional.empty:
            rows.append(
                {
                    "variant": suffix,
                    "common_centers": np.nan,
                    "strategy1": int(functional.strategy1.sum()),
                    "strategy2": int(functional.strategy2.sum()),
                    "N2_union": np.nan,
                    "CSC": np.nan,
                    "CSC_merged": np.nan,
                    "CSC_fdr": np.nan,
                }
            )
    return pd.DataFrame(rows)


def write_report(
    out: Path,
    variant_tables: dict[str, dict[str, Any]],
    luminance: dict[str, pd.DataFrame] | None,
    *,
    perms: int,
    workers: int,
) -> Path:
    lines: list[str] = []
    selected_windows = sorted({suffix.rsplit("_", 1)[0] for suffix in variant_tables})
    selected_signals = sorted({suffix.rsplit("_", 1)[1] for suffix in variant_tables})
    signal_names = {
        "lf30": "低频 1-30 Hz",
        "raw200": "HDF5 1-200 Hz",
    }
    window_text = " / ".join(selected_windows) if selected_windows else "未指定"
    signal_text = " / ".join(signal_names.get(s, s) for s in selected_signals)
    lines.append("# 颜色分析管线 v2 报告")
    lines.append("")
    lines.append(f"**分析日期**：{date.today().isoformat()}")
    lines.append("")
    lines.append("## 1. Executive summary")
    lines.append("")
    lines.append(
        f"本报告基于 7 位被试 21 个 HDF5，当前运行时间窗为 {window_text} ms，"
        f"信号变体为 {signal_text}。"
    )
    lines.append("")
    counts = _variant_counts_table(variant_tables)
    if not counts.empty:
        lines.append(df_to_markdown(counts))
    lines.append("")
    lines.append(
        "功能筛选标准（2026-08-05 起）：`strategy1` = 二因素 ANOVA（颜色 × 类别）"
        "类型 II 颜色主效应 p<0.05；`strategy2` = 任一类别 Welch t 检验 p<0.05"
        "（两组 ANOVA 等价于 t 检验）。MWU pooled（`strategy1_merged`）与四类各自"
        "显著性（`*_p_raw`）保留为信息列；FDR 仅作信息列，不作筛选门槛。"
    )
    lines.append("")
    lines.append(
        f"当前主时间窗：{window_text} ms；正式电极筛选使用 Task 1 的 ANOVA 颜色主效应。"
    )
    lines.append("")
    lines.append(
        "注意：跨通道不校正的 raw p 筛选与完全随机假设下的期望数量级一致"
        "（详见 stage01_selection 表），CSC 结论应视为探索性。"
    )
    lines.append("")
    for variant in all_variants():
        suffix = variant.suffix
        tables = variant_tables.get(suffix)
        if tables is None:
            continue
        lines.append(f"## 2. 变体 {suffix}（{variant.window_label} ms / {variant.signal_label}）")
        lines.append("")
        summary = tables.get("summary")
        if summary is not None and not summary.empty:
            lines.append("### 电极集合计数")
            lines.append("")
            lines.append(df_to_markdown(summary))
            lines.append("")
        spatial = tables.get("spatial")
        if spatial is not None and not spatial.empty:
            csc = spatial[spatial.CSC]
            if not csc.empty:
                lines.append(
                    "### CSC 电极（"
                    + ", ".join(
                        f"{s}:{c}" for s, c in zip(csc.subject, csc.channel)
                    )
                    + "）"
                )
                lines.append("")
        amp = tables.get("amp")
        if amp is not None and not amp.empty:
            anova = amp[amp.analysis == "task2_gray_fruit_anova"]
            lines.append(
                f"### 幅度统计：四水果 ANOVA {len(anova)} 通道，"
                f"red-vs-green MWU {int((amp.analysis == 'task3_red_green').sum())} 条"
            )
            lines.append("")
        decoding = tables.get("decoding")
        if decoding is not None and not decoding.empty:
            lines.append("### 频谱级解码（置换 p）")
            lines.append("")
            lines.append(
                decoding.groupby("analysis")
                .agg(
                    n_electrodes=("channel", "count"),
                    mean_accuracy=("accuracy", "mean"),
                    min_p=("p_perm", "min"),
                    n_p_lt_0_05=("p_perm", lambda x: int((x < 0.05).sum())),
                )
                .reset_index()
                .round(4)
                .pipe(df_to_markdown)
            )
            lines.append("")
    lines.append("## 3. 刺激亮度审计")
    lines.append("")
    if luminance is not None and "task1_category" in luminance:
        lines.append(df_to_markdown(luminance["task1_category"]))
        lines.append("")
        lines.append("详见 stage04_luminance/。")
    else:
        lines.append("未运行亮度审计。")
    lines.append("")
    lines.append("## 关键图")
    lines.append("")
    for stage in ("stage01_selection", "stage02_amplitude_spectral", "stage03_decoding"):
        fig_dir = out / stage / "figures"
        if not fig_dir.exists():
            continue
        for path in sorted(fig_dir.glob("*.png")):
            lines.append(f"- [{path.name}]({stage}/figures/{path.name})")
        lines.append("")
    decision_path = out / "stage05_hypotheses" / "decision_summary.md"
    if decision_path.exists():
        lines.append("## 预注册假设检验结果（SEEG）")
        lines.append("")
        lines.append(decision_path.read_text(encoding="utf-8"))
        lines.append("")
    lines.append("## 4. Limitations")
    lines.append("")
    lines.append(
        "1. 功能筛选使用 raw 双尾 p<0.05（跨通道/通道内 FDR 均不加门槛）；"
        "在该阈值下显著通道数量与全零假设期望一致，CSC 为探索性集合。"
    )
    lines.append(
        "2. 解码为频谱级（单窗 16 频带特征）探索，置换 1000 次；"
        "逐时间点解码与 cluster 校正在下一阶段实现，本报告不报告未校正的逐时间点显著窗。"
    )
    lines.append(
        "3. 规范二使用论文群体 Talairach→MNI 坐标和 20 mm 阈值，不是被试个体化 fMRI。"
    )
    lines.append(
        "4. 记忆颜色类内方向一致性诊断见 stage03_decoding/；"
        "类内方向翻转会削弱线性 cross-fruit 解码。"
    )
    lines.append(
        "5. 样本重复检查：Task 1 无重复（70 图/条件）；Task 3 每色仅 3 张唯一图片"
        "，但 HDF5 未保留 trial-level 图片身份，因此不能把平均 epoch/图片数写成精确重复次数；"
        "物理颜色解码可能部分受样本级特征影响。Task 2 每水果 15 张（约 4 次重复），"
        "但记忆颜色解码采用 leave-one-fruit-pair-out，对样本泄漏免疫。详见"
        " stage06_exploration/exemplar_identity_audit_actual_trials.csv。"
    )
    lines.append("")
    lines.append("## 5. Reproducibility")
    lines.append("")
    lines.append(
        f"运行入口：`analysis/run_final_analysis.py --perms {perms} --workers {workers}`；"
        "参数见 analysis_parameters.json。"
    )
    lines.append("")
    report = out / "report" / "final_analysis_report.md"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text("\n".join(lines), encoding="utf-8")
    return report


def write_pptx(out: Path, variant_tables: dict[str, dict[str, Any]]) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text = box.text_frame
    text.text = f"颜色分析管线 v2 ({date.today().isoformat()})"
    p = text.add_paragraph()
    selected_windows = sorted({suffix.rsplit("_", 1)[0] for suffix in variant_tables})
    selected_signals = sorted({suffix.rsplit("_", 1)[1] for suffix in variant_tables})
    p.text = f"{'/'.join(selected_windows)} ms × {'/'.join(selected_signals)} · 频谱级解码 · 亮度审计"
    p.font.size = Pt(14)
    for variant in all_variants():
        tables = variant_tables.get(variant.suffix)
        if tables is None:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(6)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.text = f"{variant.window_label} ms / {variant.signal_label}"
        summary = tables.get("summary")
        if summary is not None and not summary.empty:
            line = (
                "S1={}, S2={}, N2={}, CSC={}, CSC_fdr={}".format(
                    int(summary.strategy1.sum()),
                    int(summary.strategy2.sum()),
                    int(summary.N2_union.sum()),
                    int(summary.CSC.sum()),
                    int(
                        summary.CSC_fdr.sum()
                        if "CSC_fdr" in summary.columns
                        else 0
                    ),
                )
            )
            p = frame.add_paragraph()
            p.text = line
        decoding = tables.get("decoding")
        if decoding is not None and not decoding.empty:
            p = frame.add_paragraph()
            p.text = decoding.groupby("analysis").agg(
                n=("channel", "count"),
                mean_acc=("accuracy", "mean"),
                min_p=("p_perm", "min"),
            ).round(4).to_string()
    path = out / "report" / "final_analysis_report.pptx"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path
